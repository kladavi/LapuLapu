#!/usr/bin/env python3
"""Extract text from a PowerPoint (.pptx) file using python-pptx. Prints plain text to stdout."""
import sys
import io

# Force UTF-8 stdout on Windows
if sys.platform == "win32":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation


def extract(path: str) -> str:
    prs = Presentation(path)
    pages = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append(" | ".join(cells))
        pages.append("\n".join(lines))
    return "\n\n".join(pages)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: extract-pptx.py <path>", file=sys.stderr)
        sys.exit(1)
    try:
        print(extract(sys.argv[1]))
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
